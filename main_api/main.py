import io
import json
import os
import re
import tempfile

import httpx
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from groq import Groq
from pydantic import BaseModel, Field, field_validator
from slowapi import Limiter
from slowapi.errors import RateLimitExceeded
from slowapi.util import get_remote_address
from starlette.responses import JSONResponse

load_dotenv()

GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
IMAGE_SERVICE_URL = os.getenv("IMAGE_SERVICE_URL", "http://localhost:8001")
GROQ_MODEL = os.getenv("GROQ_MODEL", "compound-beta")

limiter = Limiter(key_func=get_remote_address)
app = FastAPI(title="Photon AI – Main API (Orchestrator)")
app.state.limiter = limiter

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.exception_handler(RateLimitExceeded)
async def rate_limit_handler(request: Request, exc: RateLimitExceeded):
    return JSONResponse(
        status_code=429,
        content={"detail": "Слишком много запросов. Попробуйте через 10 минут."},
    )


# ---------------------------------------------------------------------------
# Request / Response models
# ---------------------------------------------------------------------------

VALID_STYLES = {"formal", "creative", "pitch"}
VALID_LANGUAGES = {"ru", "kz", "en"}
VALID_CONTENT_SIZES = {"standard", "detailed"}
VALID_FORMATS = {"pdf", "pptx"}


class GenerateRequest(BaseModel):
    topic: str = Field(..., min_length=1, max_length=300)
    slides_count: int = Field(..., ge=1, le=15)
    style: str = "formal"
    language: str = "ru"
    content_size: str = "standard"
    include_images: bool = True
    format: str = "pdf"

    @field_validator("style")
    @classmethod
    def validate_style(cls, v: str) -> str:
        v = v.lower()
        if v not in VALID_STYLES:
            raise ValueError(f"style must be one of {VALID_STYLES}")
        return v

    @field_validator("language")
    @classmethod
    def validate_language(cls, v: str) -> str:
        v = v.lower()
        if v not in VALID_LANGUAGES:
            raise ValueError(f"language must be one of {VALID_LANGUAGES}")
        return v

    @field_validator("content_size")
    @classmethod
    def validate_content_size(cls, v: str) -> str:
        v = v.lower()
        if v not in VALID_CONTENT_SIZES:
            raise ValueError(f"content_size must be one of {VALID_CONTENT_SIZES}")
        return v

    @field_validator("format")
    @classmethod
    def validate_format(cls, v: str) -> str:
        v = v.lower()
        if v not in VALID_FORMATS:
            raise ValueError(f"format must be one of {VALID_FORMATS}")
        return v


# ---------------------------------------------------------------------------
# Groq helpers
# ---------------------------------------------------------------------------

LANGUAGE_NAMES = {"ru": "Russian", "kz": "Kazakh", "en": "English"}

STYLE_PROMPTS = {
    "formal": "Use an academic, professional tone. Focus on facts and structure.",
    "creative": "Use an inspiring, engaging, creative tone. Make it memorable.",
    "pitch": "Use a startup-pitch tone: identify problem, present solution, show results.",
}

BULLET_COUNTS = {
    "standard": "4-5 bullet points per slide",
    "detailed": "5-6 bullet points per slide",
}


def _groq_client() -> Groq:
    if not GROQ_API_KEY:
        raise HTTPException(status_code=500, detail="GROQ_API_KEY is not configured.")
    return Groq(api_key=GROQ_API_KEY)


def _extract_json(text: str) -> dict:
    """Extract the first JSON object from an LLM response."""
    match = re.search(r"\{[\s\S]*\}", text)
    if not match:
        raise ValueError("No JSON object found in LLM response.")
    return json.loads(match.group())


def generate_structure(client: Groq, topic: str, slides_count: int) -> dict:
    system_prompt = (
        "You are a presentation structure generator. "
        "Return ONLY a valid JSON object – no markdown, no explanation. "
        "Schema: {\"title\": string, \"slides\": [{\"type\": string, \"title\": string, \"bullets\": []}]}"
    )
    user_prompt = (
        f"Create a presentation structure with exactly {slides_count} slides on the topic: \"{topic}\". "
        "Return strictly JSON."
    )
    response = client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.4,
        max_tokens=2000,
    )
    raw = response.choices[0].message.content
    return _extract_json(raw)


def generate_content(
    client: Groq,
    structure: dict,
    topic: str,
    style: str,
    language: str,
    content_size: str,
) -> dict:
    lang_name = LANGUAGE_NAMES.get(language, "English")
    style_instruction = STYLE_PROMPTS.get(style, STYLE_PROMPTS["formal"])
    bullet_count = BULLET_COUNTS.get(content_size, BULLET_COUNTS["standard"])

    system_prompt = (
        f"You are a professional presentation writer. {style_instruction} "
        f"Write all content in {lang_name}. "
        f"Each slide must have {bullet_count} (each max 15 words). "
        "Return ONLY a valid JSON object matching the input structure – no markdown, no explanation. "
        "Schema: same as input but with bullets filled in."
    )
    user_prompt = (
        f"Fill in the following presentation structure about \"{topic}\" with real content.\n"
        f"Structure:\n{json.dumps(structure, ensure_ascii=False)}\n"
        "Return the completed JSON."
    )
    response = client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.6 if style == "creative" else 0.4,
        max_tokens=3000,
    )
    raw = response.choices[0].message.content
    result = _extract_json(raw)
    # Enforce max 6 bullets per slide
    for slide in result.get("slides", []):
        if isinstance(slide.get("bullets"), list):
            slide["bullets"] = slide["bullets"][:6]
    return result


# ---------------------------------------------------------------------------
# Image service helper
# ---------------------------------------------------------------------------

async def fetch_image(topic: str, slide_title: str, language: str) -> str | None:
    query = f"{slide_title} {topic} illustration"
    try:
        async with httpx.AsyncClient(timeout=8) as client:
            resp = await client.post(
                f"{IMAGE_SERVICE_URL}/get-image",
                json={"query": query, "language": language},
            )
            if resp.status_code == 200:
                data = resp.json()
                url = data.get("url")
                if url:
                    return url
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# Export helpers
# ---------------------------------------------------------------------------

def export_pdf(presentation: dict) -> bytes:
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=landscape(A4),
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SlideTitle",
        parent=styles["Heading1"],
        fontSize=28,
        textColor=colors.HexColor("#1c56a1"),
        spaceAfter=10,
    )
    pres_title_style = ParagraphStyle(
        "PresentationTitle",
        parent=styles["Heading1"],
        fontSize=36,
        textColor=colors.HexColor("#1c56a1"),
        alignment=1,
        spaceAfter=20,
    )
    bullet_style = ParagraphStyle(
        "Bullet",
        parent=styles["Normal"],
        fontSize=14,
        leftIndent=20,
        spaceAfter=6,
    )

    story = []
    pres_title = presentation.get("title", "Presentation")
    story.append(Paragraph(pres_title, pres_title_style))
    story.append(Spacer(1, 1 * cm))

    for i, slide in enumerate(presentation.get("slides", [])):
        if i > 0:
            story.append(PageBreak())
        story.append(Paragraph(slide.get("title", ""), title_style))
        story.append(Spacer(1, 0.3 * cm))
        for bullet in slide.get("bullets", []):
            story.append(Paragraph(f"• {bullet}", bullet_style))

    doc.build(story)
    return buf.getvalue()


def export_pptx(presentation: dict) -> bytes:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = PptxPresentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    title_layout = prs.slide_layouts[0]  # title slide

    # Title slide
    slide = prs.slides.add_slide(title_layout)
    title_ph = slide.shapes.title
    if title_ph:
        title_ph.text = presentation.get("title", "Presentation")
        title_ph.text_frame.paragraphs[0].font.size = Pt(40)
        title_ph.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1C, 0x56, 0xA1)

    # Content slides
    for slide_data in presentation.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)

        # Title box
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1C, 0x56, 0xA1)

        # Bullets box
        bullets = slide_data.get("bullets", [])
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.7), Inches(12.33), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for idx, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(18)
                p.space_after = Pt(6)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@app.get("/health")
async def health():
    return {"status": "ok", "service": "main-api"}


@app.post("/generate")
@limiter.limit("3/10minutes")
async def generate(request: Request, body: GenerateRequest):
    client = _groq_client()

    # Step 1: Generate structure
    try:
        structure = generate_structure(client, body.topic, body.slides_count)
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Structure generation failed: {e}")

    # Step 2: Generate content
    try:
        presentation = generate_content(
            client,
            structure,
            body.topic,
            body.style,
            body.language,
            body.content_size,
        )
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Content generation failed: {e}")

    # Step 3: Fetch images (optional)
    if body.include_images:
        for slide in presentation.get("slides", []):
            image_url = await fetch_image(body.topic, slide.get("title", ""), body.language)
            if image_url:
                slide["image_url"] = image_url

    # Step 4: Export
    try:
        if body.format == "pdf":
            file_bytes = export_pdf(presentation)
            media_type = "application/pdf"
            filename = "presentation.pdf"
        else:
            file_bytes = export_pptx(presentation)
            media_type = (
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            )
            filename = "presentation.pptx"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Export failed: {e}")

    return StreamingResponse(
        io.BytesIO(file_bytes),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
