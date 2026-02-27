import io
import os
import urllib.request
import textwrap
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.colors import HexColor

PDF_WIDTH, PDF_HEIGHT = 1280, 720

FONTS = {
    "Roboto-Regular": "https://github.com/googlefonts/roboto/raw/main/src/hinted/Roboto-Regular.ttf",
    "Roboto-Bold": "https://github.com/googlefonts/roboto/raw/main/src/hinted/Roboto-Bold.ttf",
    "Roboto-Black": "https://github.com/googlefonts/roboto/raw/main/src/hinted/Roboto-Black.ttf",
    "Roboto-Light": "https://github.com/googlefonts/roboto/raw/main/src/hinted/Roboto-Light.ttf"
}

for name, url in FONTS.items():
    font_path = f"{name}.ttf"
    if not os.path.exists(font_path):
        urllib.request.urlretrieve(url, font_path)
    pdfmetrics.registerFont(TTFont(name, font_path))

def draw_header(c, title, subtitle, theme, x=100, y_start=PDF_HEIGHT-100, width=1080):
    c.setFillColor(HexColor(theme["primary"]))
    c.rect(x, y_start + 5, 8, -45, fill=1, stroke=0)
    
    c.setFillColor(HexColor(theme["text_dark"]))
    c.setFont("Roboto-Black", 46)
    c.drawString(x + 25, y_start - 35, title[:50])
    
    if subtitle:
        c.setFillColor(HexColor(theme["text_gray"]))
        c.setFont("Roboto-Light", 24)
        lines = textwrap.wrap(subtitle, width=int(width/15))
        sub_y = y_start - 80
        for line in lines:
            c.drawString(x + 25, sub_y, line)
            sub_y -= 35
    return sub_y if subtitle else y_start - 80

def export_to_pdf(data):
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=(PDF_WIDTH, PDF_HEIGHT))
    theme = data.get("theme", {})
    bg_color = HexColor(theme["bg"])

    for i, slide in enumerate(data.get("slides", [])):
        layout = slide.get("layout", "grid_2x2")
        items = slide.get("items", [])
        
        # Фон
        c.setFillColor(bg_color)
        c.rect(0, 0, PDF_WIDTH, PDF_HEIGHT, fill=1, stroke=0)

        if layout == "hero_center":
            # Главный титульный экран
            c.setFillColor(HexColor(theme["primary"]))
            c.setFont("Roboto-Black", 72)
            y = PDF_HEIGHT / 2 + 50
            for line in textwrap.wrap(slide.get("title", data.get("title", "")), width=25):
                c.drawCentredString(PDF_WIDTH / 2, y, line)
                y -= 85
                
            c.setFillColor(HexColor(theme["text_gray"]))
            c.setFont("Roboto-Light", 28)
            for line in textwrap.wrap(slide.get("subtitle", ""), width=50):
                c.drawCentredString(PDF_WIDTH / 2, y - 20, line)
                y -= 40

        elif layout == "grid_3_cols":
            draw_header(c, slide.get("title", ""), slide.get("subtitle", ""), theme)
            col_w = (PDF_WIDTH - 280) / 3
            x_start = 100
            y_start = PDF_HEIGHT / 2 - 40
            
            for idx, item in enumerate(items[:3]):
                x = x_start + idx * (col_w + 40)
                
                c.setFillColor(HexColor(theme["accent"]))
                c.roundRect(x, y_start, 60, 60, 10, fill=1, stroke=0)
                c.setFillColor(HexColor(theme["primary"]))
                c.setFont("Roboto-Black", 32)
                c.drawString(x + 10, y_start + 18, f"0{idx+1}")
                
                c.setFillColor(HexColor(theme["text_dark"]))
                c.setFont("Roboto-Bold", 24)
                c.drawString(x, y_start - 40, item.get("heading", "")[:25])
                
                c.setFillColor(HexColor(theme["text_gray"]))
                c.setFont("Roboto-Regular", 18)
                text_y = y_start - 75
                for line in textwrap.wrap(item.get("text", ""), width=32):
                    c.drawString(x, text_y, line)
                    text_y -= 25

        elif layout == "grid_2x2":
            draw_header(c, slide.get("title", ""), slide.get("subtitle", ""), theme)
            col_w = (PDF_WIDTH - 300) / 2
            positions = [
                (100, PDF_HEIGHT/2 - 40), (100 + col_w + 80, PDF_HEIGHT/2 - 40),
                (100, PDF_HEIGHT/2 - 200), (100 + col_w + 80, PDF_HEIGHT/2 - 200)
            ]
            
            for idx, item in enumerate(items[:4]):
                x, y = positions[idx]
                c.setFillColor(HexColor(theme["primary"]))
                c.setFont("Roboto-Bold", 24)
                c.drawString(x, y, "• " + item.get("heading", "")[:35])
                
                c.setFillColor(HexColor(theme["text_gray"]))
                c.setFont("Roboto-Regular", 20)
                text_y = y - 35
                for line in textwrap.wrap(item.get("text", ""), width=50):
                    c.drawString(x + 20, text_y, line)
                    text_y -= 28

        elif layout == "large_cards":
            draw_header(c, slide.get("title", ""), slide.get("subtitle", ""), theme)
            y_start = PDF_HEIGHT / 2
            
            for idx, item in enumerate(items[:2]):
                x = 100
                y = y_start - (idx * 160)
                
                c.setFillColor(HexColor(theme["primary"]))
                c.setFont("Roboto-Bold", 28)
                c.drawString(x, y, item.get("heading", ""))
                
                c.setFillColor(HexColor(theme["text_gray"]))
                c.setFont("Roboto-Regular", 22)
                text_y = y - 40
                for line in textwrap.wrap(item.get("text", ""), width=80):
                    c.drawString(x, text_y, line)
                    text_y -= 30

        c.showPage()
        
    c.save()
    buffer.seek(0)
    return buffer

def export_to_pptx(data):
    from pptx.util import Inches
    from pptx import Presentation
    
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    
    for i, slide_data in enumerate(data.get("slides", [])):
        slide = prs.slides.add_slide(prs.slide_layouts[1] if i > 0 else prs.slide_layouts[0])
        slide.shapes.title.text = slide_data.get("title", data.get("title", ""))
        
        if i > 0:
            tf = slide.shapes.placeholders[1].text_frame
            p = tf.add_paragraph()
            p.text = slide_data.get("subtitle", "")
            p.font.italic = True
            
            for item in slide_data.get("items", []):
                p_head = tf.add_paragraph()
                p_head.text = f"\n{item.get('heading', '')}"
                p_head.font.bold = True
                
                p_text = tf.add_paragraph()
                p_text.text = item.get("text", "")
                p_text.level = 1
                
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer